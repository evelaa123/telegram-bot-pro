"""
Presentation generation service.
Creates PPTX files using GigaChat for content, design, and images.
GigaChat determines ALL design aspects: colors, fonts, sizes, layouts, image positions.
"""
import io
import json
import re
from typing import Optional, Dict, Any, List, Tuple
from decimal import Decimal

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from bot.services.gigachat_service import gigachat_service
from bot.services.cometapi_service import cometapi_service
from config import settings
import structlog

logger = structlog.get_logger()


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return RGBColor(51, 51, 51)  # Default gray
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        return RGBColor(51, 51, 51)


class PresentationService:
    """
    Service for generating PowerPoint presentations.
    GigaChat determines ALL design aspects including colors, fonts, sizes, and layouts.
    """
    
    def __init__(self):
        self.gigachat = gigachat_service
        self.cometapi = cometapi_service
    
    async def generate_presentation(
        self,
        topic: str,
        slides_count: int = 5,
        style: str = "business",
        include_images: bool = True,
        language: str = "ru",
        progress_callback=None
    ) -> Tuple[bytes, Dict[str, Any]]:
        """
        Generate a complete presentation with GigaChat-determined design.
        """
        slides_count = max(3, min(15, slides_count))
        total_usage = {
            "text_tokens": 0,
            "images_generated": 0,
            "total_cost_usd": Decimal("0"),
            "provider": "gigachat"
        }
        
        try:
            # Step 1: Generate full presentation specification from GigaChat
            if progress_callback:
                await progress_callback({
                    "step": "design",
                    "message": "GigaChat разрабатывает дизайн презентации..." if language == "ru" else "GigaChat designing presentation...",
                    "progress": 10
                })
            
            spec, spec_usage = await self._generate_presentation_spec(
                topic=topic,
                slides_count=slides_count,
                style_hint=style,
                include_images=include_images,
                language=language
            )
            
            total_usage["text_tokens"] += spec_usage.get("total_tokens", 0)
            if "cost_rub" in spec_usage:
                total_usage["total_cost_usd"] += spec_usage["cost_rub"] / Decimal("90")
            
            # Step 2: Create presentation from spec
            if progress_callback:
                await progress_callback({
                    "step": "creating",
                    "message": "Создание презентации..." if language == "ru" else "Creating presentation...",
                    "progress": 30
                })
            
            prs = Presentation()
            
            # Set slide dimensions from spec
            design = spec.get("design", {})
            slide_width = design.get("slide_width_inches", 13.33)
            slide_height = design.get("slide_height_inches", 7.5)
            prs.slide_width = Inches(slide_width)
            prs.slide_height = Inches(slide_height)
            
            slides_data = spec.get("slides", [])
            
            # Step 3: Generate each slide
            for i, slide_spec in enumerate(slides_data):
                progress = 30 + int((i + 1) / len(slides_data) * 60)
                
                if progress_callback:
                    await progress_callback({
                        "step": "slide",
                        "message": f"Создание слайда {i + 1}/{len(slides_data)}..." if language == "ru" else f"Creating slide {i + 1}/{len(slides_data)}...",
                        "progress": progress
                    })
                
                # Generate image if specified
                image_bytes = None
                image_spec = slide_spec.get("image")
                if include_images and image_spec and image_spec.get("prompt"):
                    image_bytes = await self._generate_image(image_spec.get("prompt"))
                    if image_bytes:
                        total_usage["images_generated"] += 1
                
                # Create slide from spec
                self._create_slide_from_spec(prs, slide_spec, design, image_bytes)
            
            # Step 4: Save presentation
            if progress_callback:
                await progress_callback({
                    "step": "saving",
                    "message": "Сохранение презентации..." if language == "ru" else "Saving presentation...",
                    "progress": 95
                })
            
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_bytes = pptx_buffer.getvalue()
            
            if progress_callback:
                await progress_callback({
                    "step": "complete",
                    "message": "Готово!" if language == "ru" else "Complete!",
                    "progress": 100
                })
            
            return pptx_bytes, {
                "title": spec.get("title", topic),
                "slides_count": len(slides_data),
                "design_style": design.get("style_name", style),
                "usage": total_usage
            }
            
        except Exception as e:
            logger.error("Presentation generation failed", error=str(e), topic=topic)
            raise
    
    async def _generate_presentation_spec(
        self,
        topic: str,
        slides_count: int,
        style_hint: str,
        include_images: bool,
        language: str
    ) -> Tuple[Dict[str, Any], Dict[str, Any]]:
        """
        Generate complete presentation specification from GigaChat.
        GigaChat determines ALL design aspects.
        """
        
        system_prompt = """Ты — профессиональный дизайнер презентаций. Твоя задача — создать ПОЛНУЮ спецификацию презентации в формате JSON.

Ты САМОСТОЯТЕЛЬНО определяешь ВСЕ аспекты дизайна на основе темы и стиля:
- Цветовую схему (подходящую к теме)
- Шрифты и их размеры
- Расположение элементов на каждом слайде
- Позиции и размеры изображений
- Декоративные элементы

Верни СТРОГО JSON в следующем формате:
{
    "title": "Название презентации",
    "design": {
        "style_name": "Название стиля (например: Корпоративный синий, Креативный градиент, Минимализм)",
        "slide_width_inches": 13.33,
        "slide_height_inches": 7.5,
        "colors": {
            "primary": "#003366",
            "secondary": "#0066CC",
            "accent": "#FF9900",
            "text": "#333333",
            "text_light": "#FFFFFF",
            "background": "#FFFFFF"
        },
        "fonts": {
            "title": {"name": "Arial", "size_pt": 44, "bold": true},
            "subtitle": {"name": "Arial", "size_pt": 24, "bold": false},
            "heading": {"name": "Arial", "size_pt": 32, "bold": true},
            "body": {"name": "Arial", "size_pt": 20, "bold": false},
            "caption": {"name": "Arial", "size_pt": 12, "bold": false}
        }
    },
    "slides": [
        {
            "slide_number": 1,
            "type": "title",
            "layout": {
                "background_color": "#003366",
                "elements": [
                    {
                        "type": "shape",
                        "shape": "rectangle",
                        "x_inches": 0,
                        "y_inches": 2.5,
                        "width_inches": 13.33,
                        "height_inches": 2.5,
                        "fill_color": "#0066CC"
                    },
                    {
                        "type": "text",
                        "content": "Заголовок презентации",
                        "x_inches": 0.5,
                        "y_inches": 2.8,
                        "width_inches": 12.33,
                        "height_inches": 1.2,
                        "font": "title",
                        "color": "#FFFFFF",
                        "align": "center"
                    },
                    {
                        "type": "text",
                        "content": "Подзаголовок",
                        "x_inches": 0.5,
                        "y_inches": 4.2,
                        "width_inches": 12.33,
                        "height_inches": 0.8,
                        "font": "subtitle",
                        "color": "#FFFFFF",
                        "align": "center"
                    }
                ]
            }
        },
        {
            "slide_number": 2,
            "type": "content",
            "title": "Заголовок слайда",
            "layout": {
                "background_color": "#FFFFFF",
                "elements": [
                    {
                        "type": "shape",
                        "shape": "rectangle",
                        "x_inches": 0,
                        "y_inches": 0,
                        "width_inches": 13.33,
                        "height_inches": 1.2,
                        "fill_color": "#003366"
                    },
                    {
                        "type": "text",
                        "content": "Заголовок слайда",
                        "x_inches": 0.5,
                        "y_inches": 0.3,
                        "width_inches": 12,
                        "height_inches": 0.8,
                        "font": "heading",
                        "color": "#FFFFFF",
                        "align": "left"
                    },
                    {
                        "type": "bullets",
                        "items": ["Пункт 1", "Пункт 2", "Пункт 3"],
                        "x_inches": 0.5,
                        "y_inches": 1.5,
                        "width_inches": 7,
                        "height_inches": 5.5,
                        "font": "body",
                        "color": "#333333",
                        "bullet_color": "#0066CC"
                    }
                ]
            },
            "image": {
                "prompt": "Описание изображения на английском для генерации",
                "x_inches": 8,
                "y_inches": 1.5,
                "width_inches": 4.8,
                "height_inches": 4.8
            }
        }
    ]
}

ПРАВИЛА:
1. Подбирай цвета, подходящие к ТЕМЕ презентации
2. Для делового стиля — сдержанные цвета (синий, серый)
3. Для креативного — яркие, контрастные
4. Для образовательного — зелёный, жёлтый, синий
5. Размещай изображения ЛОГИЧНО — справа или слева от текста
6. Первый слайд — ТИТУЛЬНЫЙ (тип "title")
7. Последний слайд — ЗАКЛЮЧИТЕЛЬНЫЙ с благодарностью
8. image.prompt должен быть на АНГЛИЙСКОМ языке
9. Используй разные layouts для разнообразия
10. Учитывай что элементы не должны перекрываться"""

        user_prompt = f"""Создай презентацию:
- Тема: {topic}
- Количество слайдов: {slides_count}
- Стиль (подсказка): {style_hint}
- Генерировать изображения: {"да" if include_images else "нет"}
- Язык контента: {"русский" if language == "ru" else "английский"}

Придумай подходящий дизайн, цвета и расположение элементов."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
        
        content, usage = await self.gigachat.generate_text(
            messages,
            model="GigaChat-2-Max",
            max_tokens=8192,
            temperature=0.8
        )
        
        # Parse JSON from response
        try:
            # Find JSON in response
            json_match = re.search(r'\{[\s\S]*\}', content)
            if json_match:
                spec = json.loads(json_match.group())
            else:
                raise ValueError("No JSON found in response")
        except (json.JSONDecodeError, ValueError) as e:
            logger.error("Failed to parse presentation spec", error=str(e), content=content[:500])
            # Return fallback spec
            spec = self._get_fallback_spec(topic, slides_count, language)
        
        return spec, usage
    
    def _get_fallback_spec(self, topic: str, slides_count: int, language: str) -> Dict[str, Any]:
        """Return fallback specification if GigaChat fails."""
        return {
            "title": topic,
            "design": {
                "style_name": "Default",
                "slide_width_inches": 13.33,
                "slide_height_inches": 7.5,
                "colors": {
                    "primary": "#003366",
                    "secondary": "#0066CC",
                    "accent": "#FF9900",
                    "text": "#333333",
                    "text_light": "#FFFFFF",
                    "background": "#FFFFFF"
                },
                "fonts": {
                    "title": {"name": "Arial", "size_pt": 44, "bold": True},
                    "subtitle": {"name": "Arial", "size_pt": 24, "bold": False},
                    "heading": {"name": "Arial", "size_pt": 32, "bold": True},
                    "body": {"name": "Arial", "size_pt": 20, "bold": False},
                    "caption": {"name": "Arial", "size_pt": 12, "bold": False}
                }
            },
            "slides": [
                {
                    "slide_number": 1,
                    "type": "title",
                    "layout": {
                        "background_color": "#003366",
                        "elements": [
                            {
                                "type": "text",
                                "content": topic,
                                "x_inches": 0.5,
                                "y_inches": 3,
                                "width_inches": 12.33,
                                "height_inches": 1.5,
                                "font": "title",
                                "color": "#FFFFFF",
                                "align": "center"
                            }
                        ]
                    }
                },
                {
                    "slide_number": slides_count,
                    "type": "end",
                    "layout": {
                        "background_color": "#003366",
                        "elements": [
                            {
                                "type": "text",
                                "content": "Спасибо за внимание!" if language == "ru" else "Thank you!",
                                "x_inches": 0.5,
                                "y_inches": 3,
                                "width_inches": 12.33,
                                "height_inches": 1.5,
                                "font": "title",
                                "color": "#FFFFFF",
                                "align": "center"
                            }
                        ]
                    }
                }
            ]
        }
    
    def _create_slide_from_spec(
        self,
        prs: Presentation,
        slide_spec: Dict[str, Any],
        design: Dict[str, Any],
        image_bytes: Optional[bytes]
    ):
        """Create a slide from GigaChat specification."""
        # Use blank layout
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        layout = slide_spec.get("layout", {})
        colors = design.get("colors", {})
        fonts = design.get("fonts", {})
        
        # Set background
        bg_color = layout.get("background_color", colors.get("background", "#FFFFFF"))
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = hex_to_rgb(bg_color)
        background.line.fill.background()
        
        # Process elements
        elements = layout.get("elements", [])
        for element in elements:
            self._add_element(slide, element, colors, fonts, prs)
        
        # Add image if provided
        image_spec = slide_spec.get("image")
        if image_bytes and image_spec:
            try:
                image_stream = io.BytesIO(image_bytes)
                x = Inches(image_spec.get("x_inches", 8))
                y = Inches(image_spec.get("y_inches", 1.5))
                width = Inches(image_spec.get("width_inches", 4.8))
                slide.shapes.add_picture(image_stream, x, y, width=width)
            except Exception as e:
                logger.warning("Failed to add image to slide", error=str(e))
    
    def _add_element(
        self,
        slide,
        element: Dict[str, Any],
        colors: Dict[str, str],
        fonts: Dict[str, Dict],
        prs: Presentation
    ):
        """Add an element to the slide based on spec."""
        element_type = element.get("type")
        
        if element_type == "shape":
            self._add_shape(slide, element, prs)
        elif element_type == "text":
            self._add_text(slide, element, colors, fonts)
        elif element_type == "bullets":
            self._add_bullets(slide, element, colors, fonts)
    
    def _add_shape(self, slide, element: Dict[str, Any], prs: Presentation):
        """Add a shape element."""
        shape_type = element.get("shape", "rectangle")
        x = Inches(element.get("x_inches", 0))
        y = Inches(element.get("y_inches", 0))
        width = Inches(element.get("width_inches", 1))
        height = Inches(element.get("height_inches", 1))
        
        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        }
        
        mso_shape = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(mso_shape, x, y, width, height)
        
        fill_color = element.get("fill_color")
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        
        # Remove border by default
        shape.line.fill.background()
    
    def _add_text(self, slide, element: Dict[str, Any], colors: Dict[str, str], fonts: Dict[str, Dict]):
        """Add a text element."""
        x = Inches(element.get("x_inches", 0))
        y = Inches(element.get("y_inches", 0))
        width = Inches(element.get("width_inches", 10))
        height = Inches(element.get("height_inches", 1))
        
        textbox = slide.shapes.add_textbox(x, y, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = str(element.get("content", ""))
        
        # Apply font settings
        font_key = element.get("font", "body")
        font_spec = fonts.get(font_key, {"name": "Arial", "size_pt": 20, "bold": False})
        
        p.font.name = font_spec.get("name", "Arial")
        p.font.size = Pt(font_spec.get("size_pt", 20))
        p.font.bold = font_spec.get("bold", False)
        
        # Apply color
        color = element.get("color", colors.get("text", "#333333"))
        p.font.color.rgb = hex_to_rgb(color)
        
        # Apply alignment
        align = element.get("align", "left")
        align_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT
        }
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    
    def _add_bullets(self, slide, element: Dict[str, Any], colors: Dict[str, str], fonts: Dict[str, Dict]):
        """Add a bullet list element."""
        x = Inches(element.get("x_inches", 0))
        y = Inches(element.get("y_inches", 0))
        width = Inches(element.get("width_inches", 10))
        height = Inches(element.get("height_inches", 5))
        
        textbox = slide.shapes.add_textbox(x, y, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        items = element.get("items", [])
        font_key = element.get("font", "body")
        font_spec = fonts.get(font_key, {"name": "Arial", "size_pt": 20, "bold": False})
        text_color = element.get("color", colors.get("text", "#333333"))
        
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {item}"
            p.font.name = font_spec.get("name", "Arial")
            p.font.size = Pt(font_spec.get("size_pt", 20))
            p.font.bold = font_spec.get("bold", False)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.space_after = Pt(12)
    
    async def _generate_image(self, prompt: str) -> Optional[bytes]:
        """Generate image using GigaChat or CometAPI as fallback."""
        try:
            # Try GigaChat first
            image_id, _ = await self.gigachat.generate_image_prompt(description=prompt)
            if image_id:
                return await self.gigachat.download_image(image_id)
        except Exception as e:
            logger.warning("GigaChat image generation failed", error=str(e))
        
        # Fallback to CometAPI
        try:
            image_url, _ = await self.cometapi.generate_image(
                prompt=prompt,
                size="1024x1024",
                style="natural"
            )
            return await self.cometapi.download_image(image_url)
        except Exception as e:
            logger.warning("CometAPI image generation also failed", error=str(e))
        
        return None
    
    async def generate_presentation_from_text(
        self,
        text: str,
        title: str = None,
        style: str = "business",
        include_images: bool = False,
        language: str = "ru",
        progress_callback=None
    ) -> Tuple[bytes, Dict[str, Any]]:
        """
        Generate presentation from existing text content.
        """
        # Combine title and text as topic for GigaChat
        topic = f"{title}: {text[:500]}" if title else text[:500]
        
        return await self.generate_presentation(
            topic=topic,
            slides_count=max(3, min(10, len(text) // 500 + 2)),  # Estimate slides from text length
            style=style,
            include_images=include_images,
            language=language,
            progress_callback=progress_callback
        )


# Global service instance
presentation_service = PresentationService()
