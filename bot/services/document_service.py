"""
Document processing service.
Handles parsing and extraction from various document formats.
"""
import io
from typing import Optional, Tuple, List, Dict, Any
from pathlib import Path
import structlog

logger = structlog.get_logger()


class DocumentService:
    """
    Service for processing various document formats.
    Extracts text and images for AI analysis.
    """
    
    SUPPORTED_EXTENSIONS = {
        # Text/Markup
        'txt', 'md', 'csv', 'json', 'xml',
        # PDF
        'pdf',
        # Microsoft Office
        'docx', 'xlsx', 'pptx',
        # Images (for direct vision analysis)
        'jpg', 'jpeg', 'png', 'webp', 'gif'
    }
    
    MAX_TEXT_LENGTH = 100000  # Characters
    
    def is_supported(self, filename: str) -> bool:
        """Check if file format is supported."""
        ext = Path(filename).suffix.lower().lstrip('.')
        return ext in self.SUPPORTED_EXTENSIONS
    
    def get_extension(self, filename: str) -> str:
        """Get file extension."""
        return Path(filename).suffix.lower().lstrip('.')
    
    async def process_document(
        self,
        file_data: bytes,
        filename: str,
        max_pages: int = 50,
        max_rows: int = 5000,
        max_slides: int = 100
    ) -> Tuple[str, Dict[str, Any], List[bytes]]:
        """
        Process document and extract content.
        
        Args:
            file_data: Raw file bytes
            filename: Original filename
            max_pages: Maximum PDF pages to process
            max_rows: Maximum Excel rows to process
            max_slides: Maximum PowerPoint slides to process
            
        Returns:
            Tuple of (extracted_text, metadata, extracted_images)
        """
        ext = self.get_extension(filename)
        
        processors = {
            'txt': self._process_text,
            'md': self._process_text,
            'csv': self._process_csv,
            'json': self._process_json,
            'xml': self._process_text,
            'pdf': lambda d, f: self._process_pdf(d, f, max_pages),
            'docx': self._process_docx,
            'xlsx': lambda d, f: self._process_xlsx(d, f, max_rows),
            'pptx': lambda d, f: self._process_pptx(d, f, max_slides),
        }
        
        # Image files - return empty text, metadata, and image
        if ext in {'jpg', 'jpeg', 'png', 'webp', 'gif'}:
            return "", {"type": "image", "filename": filename}, [file_data]
        
        processor = processors.get(ext)
        if not processor:
            raise ValueError(f"Unsupported file format: {ext}")
        
        try:
            return await processor(file_data, filename)
        except Exception as e:
            logger.error(f"Document processing error", filename=filename, error=str(e))
            raise
    
    async def _process_text(
        self,
        file_data: bytes,
        filename: str
    ) -> Tuple[str, Dict[str, Any], List[bytes]]:
        """Process plain text files."""
        # Try different encodings
        encodings = ['utf-8', 'cp1251', 'latin-1']
        text = None
        
        for encoding in encodings:
            try:
                text = file_data.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        
        if text is None:
            text = file_data.decode('utf-8', errors='replace')
        
        # Truncate if needed
        if len(text) > self.MAX_TEXT_LENGTH:
            text = text[:self.MAX_TEXT_LENGTH] + "\n\n[... текст обрезан ...]"
        
        metadata = {
            "type": "text",
            "filename": filename,
            "characters": len(text),
            "lines": text.count('\n') + 1
        }
        
        return text, metadata, []
    
    async def _process_csv(
        self,
        file_data: bytes,
        filename: str
    ) -> Tuple[str, Dict[str, Any], List[bytes]]:
        """Process CSV files."""
        import csv
        
        # Decode
        try:
            text = file_data.decode('utf-8')
        except UnicodeDecodeError:
            text = file_data.decode('cp1251', errors='replace')
        
        # Parse CSV
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        
        # Convert to markdown table
        if rows:
            header = rows[0]
            md_lines = [
                "| " + " | ".join(header) + " |",
                "| " + " | ".join(["---"] * len(header)) + " |"
            ]
            for row in rows[1:]:
                # Pad row to match header length
                padded = row + [''] * (len(header) - len(row))
                md_lines.append("| " + " | ".join(padded[:len(header)]) + " |")
            
            markdown_text = "\n".join(md_lines)
        else:
            markdown_text = text
        
        # Truncate if needed
        if len(markdown_text) > self.MAX_TEXT_LENGTH:
            markdown_text = markdown_text[:self.MAX_TEXT_LENGTH] + "\n\n[... данные обрезаны ...]"
        
        metadata = {
            "type": "csv",
            "filename": filename,
            "rows": len(rows),
            "columns": len(rows[0]) if rows else 0
        }
        
        return markdown_text, metadata, []
    
    async def _process_json(
        self,
        file_data: bytes,
        filename: str
    ) -> Tuple[str, Dict[str, Any], List[bytes]]:
        """Process JSON files."""
        import json
        
        text = file_data.decode('utf-8')
        
        # Pretty print JSON
        try:
            data = json.loads(text)
            pretty_text = json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            pretty_text = text
        
        # Truncate if needed
        if len(pretty_text) > self.MAX_TEXT_LENGTH:
            pretty_text = pretty_text[:self.MAX_TEXT_LENGTH] + "\n\n[... данные обрезаны ...]"
        
        metadata = {
            "type": "json",
            "filename": filename,
            "characters": len(pretty_text)
        }
        
        return pretty_text, metadata, []
    
    async def _process_pdf(
        self,
        file_data: bytes,
        filename: str,
        max_pages: int = 50
    ) -> Tuple[str, Dict[str, Any], List[bytes]]:
        """
        Process PDF files.
        First tries text extraction, falls back to image conversion for scanned PDFs.
        """
        import pdfplumber
        
        images = []
        text_parts = []
        
        with pdfplumber.open(io.BytesIO(file_data)) as pdf:
            total_pages = len(pdf.pages)
            pages_to_process = min(total_pages, max_pages)
            
            for i, page in enumerate(pdf.pages[:pages_to_process]):
                # Extract text
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(f"=== Страница {i + 1} ===\n{page_text}")
                else:
                    # Page might be scanned - convert to image
                    try:
                        img = page.to_image(resolution=150)
                        img_bytes = io.BytesIO()
                        img.save(img_bytes, format='PNG')
                        images.append(img_bytes.getvalue())
                        text_parts.append(f"=== Страница {i + 1} (изображение) ===")
                    except Exception as e:
                        logger.warning(f"Failed to convert PDF page to image", page=i, error=str(e))
                        text_parts.append(f"=== Страница {i + 1} (не удалось обработать) ===")
        
        full_text = "\n\n".join(text_parts)
        
        # Truncate if needed
        if len(full_text) > self.MAX_TEXT_LENGTH:
            full_text = full_text[:self.MAX_TEXT_LENGTH] + "\n\n[... текст обрезан ...]"
        
        metadata = {
            "type": "pdf",
            "filename": filename,
            "total_pages": total_pages,
            "processed_pages": pages_to_process,
            "image_pages": len(images),
            "has_text": bool(text_parts)
        }
        
        if total_pages > max_pages:
            metadata["warning"] = f"Обработаны только первые {max_pages} страниц из {total_pages}"
        
        return full_text, metadata, images
    
    async def _process_docx(
        self,
        file_data: bytes,
        filename: str
    ) -> Tuple[str, Dict[str, Any], List[bytes]]:
        """Process Microsoft Word documents."""
        from docx import Document
        from docx.opc.exceptions import PackageNotFoundError
        
        try:
            doc = Document(io.BytesIO(file_data))
        except PackageNotFoundError:
            raise ValueError("Invalid DOCX file format")
        
        text_parts = []
        images = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Handle headings
                if para.style.name.startswith('Heading'):
                    level = para.style.name[-1] if para.style.name[-1].isdigit() else '1'
                    text_parts.append(f"{'#' * int(level)} {para.text}")
                else:
                    text_parts.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            
            if rows:
                # Convert to markdown table
                header = rows[0]
                md_table = [
                    "| " + " | ".join(header) + " |",
                    "| " + " | ".join(["---"] * len(header)) + " |"
                ]
                for row in rows[1:]:
                    md_table.append("| " + " | ".join(row) + " |")
                text_parts.append("\n".join(md_table))
        
        # Extract images
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    images.append(rel.target_part.blob)
                except Exception:
                    pass
        
        full_text = "\n\n".join(text_parts)
        
        # Truncate if needed
        if len(full_text) > self.MAX_TEXT_LENGTH:
            full_text = full_text[:self.MAX_TEXT_LENGTH] + "\n\n[... текст обрезан ...]"
        
        metadata = {
            "type": "docx",
            "filename": filename,
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
            "images": len(images)
        }
        
        return full_text, metadata, images
    
    async def _process_xlsx(
        self,
        file_data: bytes,
        filename: str,
        max_rows: int = 5000
    ) -> Tuple[str, Dict[str, Any], List[bytes]]:
        """Process Microsoft Excel documents."""
        from openpyxl import load_workbook
        
        wb = load_workbook(io.BytesIO(file_data), data_only=True)
        
        text_parts = []
        total_rows = 0
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"## Лист: {sheet_name}\n")
            
            # Get data as list of rows
            rows = []
            for row in sheet.iter_rows(max_row=max_rows):
                cells = []
                for cell in row:
                    value = cell.value
                    if value is None:
                        cells.append("")
                    elif isinstance(value, (int, float)):
                        cells.append(str(value))
                    else:
                        cells.append(str(value).strip())
                
                # Skip completely empty rows
                if any(cells):
                    rows.append(cells)
            
            total_rows += len(rows)
            
            if rows:
                # Find max columns
                max_cols = max(len(row) for row in rows)
                
                # Normalize row lengths
                normalized_rows = [
                    row + [''] * (max_cols - len(row))
                    for row in rows
                ]
                
                # Convert to markdown table
                header = normalized_rows[0]
                md_table = [
                    "| " + " | ".join(header) + " |",
                    "| " + " | ".join(["---"] * len(header)) + " |"
                ]
                for row in normalized_rows[1:]:
                    md_table.append("| " + " | ".join(row) + " |")
                
                text_parts.append("\n".join(md_table))
        
        full_text = "\n\n".join(text_parts)
        
        # Truncate if needed
        if len(full_text) > self.MAX_TEXT_LENGTH:
            full_text = full_text[:self.MAX_TEXT_LENGTH] + "\n\n[... данные обрезаны ...]"
        
        metadata = {
            "type": "xlsx",
            "filename": filename,
            "sheets": len(wb.sheetnames),
            "sheet_names": wb.sheetnames,
            "total_rows": total_rows
        }
        
        if total_rows > max_rows:
            metadata["warning"] = f"Обработаны только первые {max_rows} строк"
        
        return full_text, metadata, []
    
    async def _process_pptx(
        self,
        file_data: bytes,
        filename: str,
        max_slides: int = 100
    ) -> Tuple[str, Dict[str, Any], List[bytes]]:
        """Process Microsoft PowerPoint documents."""
        from pptx import Presentation
        
        prs = Presentation(io.BytesIO(file_data))
        
        text_parts = []
        images = []
        
        total_slides = len(prs.slides)
        slides_to_process = min(total_slides, max_slides)
        
        for i, slide in enumerate(prs.slides[:slides_to_process]):
            slide_text = [f"## Слайд {i + 1}"]
            
            for shape in slide.shapes:
                # Extract text
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text.append(text)
                
                # Extract images
                if hasattr(shape, 'image'):
                    try:
                        images.append(shape.image.blob)
                    except Exception:
                        pass
            
            if len(slide_text) > 1:  # More than just the header
                text_parts.append("\n".join(slide_text))
        
        full_text = "\n\n".join(text_parts)
        
        # Truncate if needed
        if len(full_text) > self.MAX_TEXT_LENGTH:
            full_text = full_text[:self.MAX_TEXT_LENGTH] + "\n\n[... текст обрезан ...]"
        
        metadata = {
            "type": "pptx",
            "filename": filename,
            "total_slides": total_slides,
            "processed_slides": slides_to_process,
            "images": len(images)
        }
        
        if total_slides > max_slides:
            metadata["warning"] = f"Обработаны только первые {max_slides} слайдов из {total_slides}"
        
        return full_text, metadata, images


# Global service instance
document_service = DocumentService()
