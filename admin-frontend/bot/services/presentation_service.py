"""
Presentation generation service.
Creates PPTX files using GigaChat for content and CometAPI for images.
"""
import io
import asyncio
from typing import Optional, Dict, Any, List, Tuple
from decimal import Decimal
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from bot.services.gigachat_service import gigachat_service
from bot.services.cometapi_service import cometapi_service
from config import settings
import structlog

logger = structlog.get_logger()


class PresentationService:
    """
    Service for generating PowerPoint presentations.
    Uses GigaChat for content structure and CometAPI for images.
    """
    
    # Color schemes for different styles
    COLOR_SCHEMES = {
        "business": {
            "primary": RGBColor(0, 51, 102),      # Dark blue
            "secondary": RGBColor(0, 102, 153),   # Medium blue
            "accent": RGBColor(255, 153, 0),      # Orange
            "text": RGBColor(51, 51, 51),         # Dark gray
            "background": RGBColor(255, 255, 255) # White
        },
        "creative": {
            "primary": RGBColor(102, 45, 145),    # Purple
            "secondary": RGBColor(255, 102, 102), # Coral
            "accent": RGBColor(0, 204, 153),      # Teal
            "text": RGBColor(51, 51, 51),
            "background": RGBColor(255, 255, 255)
        },
        "educational": {
            "primary": RGBColor(0, 128, 0),       # Green
            "secondary": RGBColor(0, 102, 204),   # Blue
            "accent": RGBColor(255, 204, 0),      # Yellow
            "text": RGBColor(51, 51, 51),
            "background": RGBColor(255, 255, 255)
        },
        "modern": {
            "primary": RGBColor(33, 33, 33),      # Dark
            "secondary": RGBColor(76, 175, 80),   # Green
            "accent": RGBColor(255, 193, 7),      # Amber
            "text": RGBColor(33, 33, 33),
            "background": RGBColor(250, 250, 250)
        }
    }
    
    def __init__(self):
        self.gigachat = gigachat_service
        self.cometapi = cometapi_service
    
    async def generate_presentation(
        self,
        topic: str,
        slides_count: int = 5,
        style: str = "business",
        include_images: bool = True,
        language: str = "ru",
        progress_callback=None
    ) -> Tuple[bytes, Dict[str, Any]]:
        """
        Generate a complete presentation.
        
        Args:
            topic: Presentation topic
            slides_count: Number of slides (3-15)
            style: Visual style (business, creative, educational, modern)
            include_images: Whether to generate images for slides
            language: Output language
            progress_callback: Async callback for progress updates
            
        Returns:
            Tuple of (pptx_bytes, generation_info)
        """
        slides_count = max(3, min(15, slides_count))
        total_usage = {
            "text_tokens": 0,
            "images_generated": 0,
            "total_cost_usd": Decimal("0"),
            "provider": "gigachat+cometapi"
        }
        
        try:
            # Step 1: Generate structure with GigaChat
            if progress_callback:
                await progress_callback({
                    "step": "structure",
                    "message": "Генерация структуры презентации..." if language == "ru" else "Generating presentation structure...",
                    "progress": 10
                })
            
            structure, structure_usage = await self.gigachat.generate_presentation_structure(
                topic=topic,
                slides_count=slides_count,
                style=style,
                language=language
            )
            
            total_usage["text_tokens"] += structure_usage.get("total_tokens", 0)
            if "cost_rub" in structure_usage:
                # Convert RUB to USD (approximate)
                total_usage["total_cost_usd"] += structure_usage["cost_rub"] / Decimal("90")
            
            # Step 2: Create presentation
            if progress_callback:
                await progress_callback({
                    "step": "creating",
                    "message": "Создание презентации..." if language == "ru" else "Creating presentation...",
                    "progress": 30
                })
            
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            colors = self.COLOR_SCHEMES.get(style, self.COLOR_SCHEMES["business"])
            slides_data = structure.get("slides", [])
            
            # Step 3: Generate slides with images
            for i, slide_data in enumerate(slides_data):
                progress = 30 + int((i + 1) / len(slides_data) * 60)
                
                if progress_callback:
                    await progress_callback({
                        "step": "slide",
                        "message": f"Создание слайда {i + 1}/{len(slides_data)}..." if language == "ru" else f"Creating slide {i + 1}/{len(slides_data)}...",
                        "progress": progress
                    })
                
                # Generate image for slide if needed
                image_bytes = None
                if include_images and slide_data.get("image_prompt") and i > 0 and i < len(slides_data) - 1:
                    try:
                        image_url, image_usage = await self.cometapi.generate_image(
                            prompt=slide_data["image_prompt"],
                            size="1024x1024",
                            style="natural"
                        )
                        image_bytes = await self.cometapi.download_image(image_url)
                        total_usage["images_generated"] += 1
                        total_usage["total_cost_usd"] += image_usage.get("cost_usd", Decimal("0"))
                    except Exception as e:
                        logger.warning("Failed to generate image for slide", slide=i, error=str(e))
                
                # Create slide
                self._create_slide(
                    prs,
                    slide_data,
                    colors,
                    image_bytes,
                    is_title_slide=(i == 0),
                    is_last_slide=(i == len(slides_data) - 1)
                )
            
            # Step 4: Save presentation
            if progress_callback:
                await progress_callback({
                    "step": "saving",
                    "message": "Сохранение презентации..." if language == "ru" else "Saving presentation...",
                    "progress": 95
                })
            
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_bytes = pptx_buffer.getvalue()
            
            if progress_callback:
                await progress_callback({
                    "step": "complete",
                    "message": "Готово!" if language == "ru" else "Complete!",
                    "progress": 100
                })
            
            return pptx_bytes, {
                "title": structure.get("title", topic),
                "slides_count": len(slides_data),
                "usage": total_usage
            }
            
        except Exception as e:
            logger.error("Presentation generation failed", error=str(e), topic=topic)
            raise
    
    def _create_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        colors: Dict[str, RGBColor],
        image_bytes: Optional[bytes],
        is_title_slide: bool,
        is_last_slide: bool
    ):
        """Create a single slide in the presentation."""
        # Use blank layout
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = colors["background"]
        background.line.fill.background()
        
        if is_title_slide:
            self._create_title_slide(slide, slide_data, colors, prs)
        elif is_last_slide:
            self._create_end_slide(slide, slide_data, colors, prs)
        else:
            self._create_content_slide(slide, slide_data, colors, prs, image_bytes)
    
    def _create_title_slide(self, slide, slide_data, colors, prs):
        """Create title slide."""
        # Add decorative shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(2.5),
            prs.slide_width, Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8),
            Inches(12.33), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = slide_data.get("title", "Презентация")
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle if available
        if slide_data.get("content"):
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5),
                Inches(12.33), Inches(0.8)
            )
            sub_frame = subtitle_box.text_frame
            sub_frame.paragraphs[0].text = slide_data["content"][0] if isinstance(slide_data["content"], list) else str(slide_data["content"])
            sub_frame.paragraphs[0].font.size = Pt(24)
            sub_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(self, slide, slide_data, colors, prs, image_bytes):
        """Create content slide with optional image."""
        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = colors["primary"]
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = slide_data.get("title", "")
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Content area dimensions depend on image presence
        if image_bytes:
            content_width = Inches(7)
            content_left = Inches(0.5)
            
            # Add image on the right
            try:
                image_stream = io.BytesIO(image_bytes)
                slide.shapes.add_picture(
                    image_stream,
                    Inches(8), Inches(1.8),
                    width=Inches(4.8)
                )
            except Exception as e:
                logger.warning("Failed to add image to slide", error=str(e))
                content_width = Inches(12)
        else:
            content_width = Inches(12)
            content_left = Inches(0.5)
        
        # Content bullets
        content = slide_data.get("content", [])
        if content:
            content_box = slide.shapes.add_textbox(
                content_left, Inches(1.8),
                content_width, Inches(5)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for i, item in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"• {item}"
                p.font.size = Pt(20)
                p.font.color.rgb = colors["text"]
                p.space_after = Pt(12)
        
        # Slide number
        num_box = slide.shapes.add_textbox(
            Inches(12.5), Inches(7),
            Inches(0.5), Inches(0.3)
        )
        num_frame = num_box.text_frame
        num_frame.paragraphs[0].text = str(slide_data.get("slide_number", ""))
        num_frame.paragraphs[0].font.size = Pt(12)
        num_frame.paragraphs[0].font.color.rgb = colors["secondary"]
        num_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    def _create_end_slide(self, slide, slide_data, colors, prs):
        """Create ending slide (Thank you / Questions)."""
        # Full background color
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = colors["primary"]
        bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.33), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = slide_data.get("title", "Спасибо за внимание!")
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if slide_data.get("content"):
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2),
                Inches(12.33), Inches(1)
            )
            sub_frame = sub_box.text_frame
            sub_text = slide_data["content"][0] if isinstance(slide_data["content"], list) else str(slide_data["content"])
            sub_frame.paragraphs[0].text = sub_text
            sub_frame.paragraphs[0].font.size = Pt(24)
            sub_frame.paragraphs[0].font.color.rgb = colors["accent"]
            sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    async def generate_presentation_from_text(
        self,
        text: str,
        title: str = None,
        style: str = "business",
        include_images: bool = False,
        language: str = "ru",
        progress_callback=None
    ) -> Tuple[bytes, Dict[str, Any]]:
        """
        Generate presentation from existing text content.
        Useful when user provides their own content.
        
        Args:
            text: Text content to convert to presentation
            title: Optional presentation title
            style: Visual style
            include_images: Whether to generate images
            language: Output language
            progress_callback: Progress callback
            
        Returns:
            Tuple of (pptx_bytes, generation_info)
        """
        # Use GigaChat to structure the text into slides
        system_prompt = """Преобразуй текст в структуру презентации.

Верни JSON:
{
    "title": "Название презентации",
    "slides": [
        {
            "slide_number": 1,
            "title": "Заголовок слайда",
            "content": ["Пункт 1", "Пункт 2"],
            "image_prompt": "Image description in English"
        }
    ]
}

Правила:
- Первый слайд - титульный
- Последний слайд - заключение
- 3-5 пунктов на слайд
- Сохрани ключевые идеи текста"""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Текст:\n{text}\n\nНазвание: {title or 'Презентация'}"}
        ]
        
        content, usage = await self.gigachat.generate_text(
            messages,
            model="GigaChat",
            max_tokens=4096
        )
        
        import json
        try:
            json_start = content.find('{')
            json_end = content.rfind('}') + 1
            structure = json.loads(content[json_start:json_end])
        except:
            structure = {
                "title": title or "Презентация",
                "slides": [
                    {"slide_number": 1, "title": title or "Презентация", "content": []},
                    {"slide_number": 2, "title": "Содержание", "content": text.split('\n')[:5]},
                    {"slide_number": 3, "title": "Спасибо!", "content": ["Вопросы?"]}
                ]
            }
        
        # Now create the presentation
        total_usage = {"text_tokens": usage.get("total_tokens", 0), "images_generated": 0, "total_cost_usd": Decimal("0")}
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        colors = self.COLOR_SCHEMES.get(style, self.COLOR_SCHEMES["business"])
        
        slides_data = structure.get("slides", [])
        for i, slide_data in enumerate(slides_data):
            image_bytes = None
            if include_images and slide_data.get("image_prompt") and i > 0 and i < len(slides_data) - 1:
                try:
                    image_url, img_usage = await self.cometapi.generate_image(slide_data["image_prompt"])
                    image_bytes = await self.cometapi.download_image(image_url)
                    total_usage["images_generated"] += 1
                    total_usage["total_cost_usd"] += img_usage.get("cost_usd", Decimal("0"))
                except:
                    pass
            
            self._create_slide(prs, slide_data, colors, image_bytes, i == 0, i == len(slides_data) - 1)
        
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        
        return pptx_buffer.getvalue(), {
            "title": structure.get("title", title),
            "slides_count": len(slides_data),
            "usage": total_usage
        }


# Global service instance
presentation_service = PresentationService()
